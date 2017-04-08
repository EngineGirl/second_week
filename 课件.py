from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_slide_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "贝叶斯模型"
subtitle.text = "数据分析基础"

# 第二页
title_slide_layout2 = prs.slide_layouts[0]
slide2 = prs.slides.add_slide(title_slide_layout2)
title2 = slide2.shapes.title
subtitle2 = slide2.placeholders[1]
title2.text = "问题"
subtitle2.text = "若某种疾病在所有人中患病率为1%，若患病，检验结果为" + \
                 "阳性的概率为90%，未患病，而检验结果为阳性的概率为9%。" + \
                 "如果我被检验为阳性，那么我真正患病的概率为多少？"

# 第三页
title_slide_layout3 = prs.slide_layouts[0]
slide3 = prs.slides.add_slide(title_slide_layout3)
title3 = slide3.shapes.title
subtitle3 = slide3.placeholders[1]
title3.text = "答案"
subtitle3.text = "9%, https://github.com/Windsooon/What-Bayes-tell-us-/" + \
                 "blob/master/README_CHINESE.md"


# 第四页
img_path = "bayes.png"

blank_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(blank_slide_layout)
shapes = slide.shapes
body_shape = shapes.placeholders[1]
tf = body_shape.text_frame

left = Inches(0)
top = Inches(5)
height = Inches(1.5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)
title_shape = shapes.title
title_shape.text = '贝叶斯公式推导'
p = tf.add_paragraph()
p.text = 'http://www.sci.utah.edu/~gerig/CS6640-F2010/prob-tut.pdf'
p.level = 1

# 第五页
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Wordpress入门"
subtitle.text = "https://github.com/EngineGirl/basic-tutorial/" + \
                "blob/master/WordPress%E5%85%A5%E9%97%A8.md"


prs.save('second_week.pptx')
